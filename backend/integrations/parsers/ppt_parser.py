"""Parser for PPT/PPTX documents"""
from pathlib import Path
import logging

logger = logging.getLogger(__name__)


class PPTParser:
    """Parser for PowerPoint files using python-pptx"""

    def __init__(self):
        self.supported_extensions = [".ppt", ".pptx"]

    def parse(self, filepath):
        """Extract text from slides"""
        filepath = Path(filepath)
        if not filepath.exists():
            raise FileNotFoundError(f"PPT file not found: {filepath}")

        if filepath.suffix.lower() not in self.supported_extensions:
            raise ValueError(f"Not a PPT/PPTX file: {filepath}")

        try:
            from pptx import Presentation
        except ImportError:
            raise ImportError("python-pptx not installed. Install with: pip install python-pptx")

        prs = Presentation(filepath)
        full_text = ""
        slides = []

        for slide_idx, slide in enumerate(prs.slides, start=1):
            slide_text_parts = []
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text:
                    slide_text_parts.append(shape.text)
            slide_text = "\n".join(slide_text_parts).strip()
            if slide_text:
                slides.append({"slide_number": slide_idx, "text": slide_text})
                full_text += f"\n\n--- Slide {slide_idx} ---\n\n{slide_text}"

        return {
            "text": full_text.strip(),
            "pages": slides,  # align with pdf parser format
            "total_pages": len(slides)
        }

